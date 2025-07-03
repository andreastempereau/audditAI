"""Celery worker for async document processing."""

import asyncio
import json
import logging
import re
from typing import List, Dict, Any, Optional
from uuid import UUID
from pathlib import Path
import tempfile
import os

import redis
from sqlalchemy.ext.asyncio import create_async_engine, AsyncSession
from sqlalchemy import text
import pandas as pd

# Text extraction libraries
try:
    import tika
    from tika import parser as tika_parser
    TIKA_AVAILABLE = True
except ImportError:
    TIKA_AVAILABLE = False

try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

from app.core.config import get_settings
from app.services.data_room import StorageService
from app.workers.embedder_client import EmbedderClient

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

settings = get_settings()


class TextExtractor:
    """Extract text from various file formats."""
    
    def __init__(self):
        self.storage = StorageService()
        
    async def extract_text(self, file_data: bytes, content_type: str, filename: str) -> Dict[str, Any]:
        """Extract text from file based on content type."""
        try:
            if content_type == 'application/pdf':
                return await self._extract_pdf_text(file_data)
            elif content_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
                return await self._extract_docx_text(file_data)
            elif content_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
                return await self._extract_pptx_text(file_data)
            elif content_type == 'text/plain':
                return await self._extract_txt_text(file_data)
            elif content_type == 'text/csv':
                return await self._extract_csv_text(file_data, filename)
            elif content_type == 'application/zip':
                return await self._extract_zip_text(file_data)
            else:
                raise ValueError(f"Unsupported content type: {content_type}")
                
        except Exception as e:
            logger.error(f"Text extraction failed for {content_type}: {str(e)}")
            return {"text": "", "pages": [], "metadata": {"error": str(e)}}
    
    async def _extract_pdf_text(self, file_data: bytes) -> Dict[str, Any]:
        """Extract text from PDF using Apache Tika."""
        if not TIKA_AVAILABLE:
            raise ImportError("Apache Tika not available for PDF processing")
        
        # Write to temporary file for Tika processing
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as temp_file:
            temp_file.write(file_data)
            temp_path = temp_file.name
        
        try:
            # Extract text using Tika
            parsed = tika_parser.from_file(temp_path)
            full_text = parsed.get('content', '') or ''
            
            # Split into pages (rough approximation)
            pages = self._split_into_pages(full_text)
            
            return {
                "text": full_text,
                "pages": pages,
                "metadata": {
                    "page_count": len(pages),
                    "extractor": "tika"
                }
            }
        finally:
            # Clean up temp file
            os.unlink(temp_path)
    
    async def _extract_docx_text(self, file_data: bytes) -> Dict[str, Any]:
        """Extract text from DOCX using python-docx."""
        if not DOCX_AVAILABLE:
            raise ImportError("python-docx not available for DOCX processing")
        
        # Write to temporary file
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as temp_file:
            temp_file.write(file_data)
            temp_path = temp_file.name
        
        try:
            doc = DocxDocument(temp_path)
            
            # Extract text from paragraphs
            paragraphs = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    paragraphs.append(paragraph.text)
            
            full_text = '\n\n'.join(paragraphs)
            
            # Split into pages (rough approximation based on content length)
            pages = self._split_into_pages(full_text, chars_per_page=3000)
            
            return {
                "text": full_text,
                "pages": pages,
                "metadata": {
                    "paragraph_count": len(paragraphs),
                    "page_count": len(pages),
                    "extractor": "python-docx"
                }
            }
        finally:
            os.unlink(temp_path)
    
    async def _extract_pptx_text(self, file_data: bytes) -> Dict[str, Any]:
        """Extract text from PPTX using python-pptx."""
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx not available for PPTX processing")
        
        # Write to temporary file
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_file:
            temp_file.write(file_data)
            temp_path = temp_file.name
        
        try:
            prs = Presentation(temp_path)
            
            slides_text = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                if slide_text:
                    slides_text.append({
                        "page": slide_num,
                        "text": '\n'.join(slide_text)
                    })
            
            full_text = '\n\n'.join([slide['text'] for slide in slides_text])
            
            return {
                "text": full_text,
                "pages": slides_text,
                "metadata": {
                    "slide_count": len(slides_text),
                    "page_count": len(slides_text),
                    "extractor": "python-pptx"
                }
            }
        finally:
            os.unlink(temp_path)
    
    async def _extract_txt_text(self, file_data: bytes) -> Dict[str, Any]:
        """Extract text from plain text file."""
        try:
            # Try UTF-8 first, then fallback to latin-1
            try:
                text = file_data.decode('utf-8')
            except UnicodeDecodeError:
                text = file_data.decode('latin-1', errors='ignore')
            
            # Split into pages (rough approximation)
            pages = self._split_into_pages(text)
            
            return {
                "text": text,
                "pages": pages,
                "metadata": {
                    "page_count": len(pages),
                    "extractor": "native"
                }
            }
        except Exception as e:
            logger.error(f"Failed to extract text from TXT: {str(e)}")
            return {"text": "", "pages": [], "metadata": {"error": str(e)}}
    
    async def _extract_csv_text(self, file_data: bytes, filename: str) -> Dict[str, Any]:
        """Extract text from CSV using pandas."""
        try:
            # Write to temporary file
            with tempfile.NamedTemporaryFile(suffix='.csv', delete=False) as temp_file:
                temp_file.write(file_data)
                temp_path = temp_file.name
            
            try:
                # Read CSV
                df = pd.read_csv(temp_path)
                
                # Convert to text representation
                text_parts = [f"CSV File: {filename}"]
                text_parts.append(f"Columns: {', '.join(df.columns.tolist())}")
                text_parts.append(f"Rows: {len(df)}")
                text_parts.append("\nData Preview:")
                
                # Add sample data (first 10 rows)
                for idx, row in df.head(10).iterrows():
                    row_text = " | ".join([f"{col}: {val}" for col, val in row.items() if pd.notna(val)])
                    text_parts.append(f"Row {idx + 1}: {row_text}")
                
                full_text = '\n'.join(text_parts)
                
                # Split into pages
                pages = self._split_into_pages(full_text)
                
                return {
                    "text": full_text,
                    "pages": pages,
                    "metadata": {
                        "row_count": len(df),
                        "column_count": len(df.columns),
                        "columns": df.columns.tolist(),
                        "page_count": len(pages),
                        "extractor": "pandas"
                    }
                }
            finally:
                os.unlink(temp_path)
                
        except Exception as e:
            logger.error(f"Failed to extract text from CSV: {str(e)}")
            return {"text": "", "pages": [], "metadata": {"error": str(e)}}
    
    async def _extract_zip_text(self, file_data: bytes) -> Dict[str, Any]:
        """Extract text from ZIP archive (extract text from contained files)."""
        import zipfile
        import io
        
        try:
            text_parts = []
            pages = []
            file_count = 0
            
            with zipfile.ZipFile(io.BytesIO(file_data), 'r') as zip_ref:
                for file_info in zip_ref.filelist:
                    if file_info.is_dir():
                        continue
                    
                    filename = file_info.filename
                    file_content = zip_ref.read(filename)
                    
                    # Try to extract text from known file types
                    if filename.lower().endswith(('.txt', '.md', '.py', '.js', '.html', '.css')):
                        try:
                            file_text = file_content.decode('utf-8', errors='ignore')
                            text_parts.append(f"\n--- File: {filename} ---\n")
                            text_parts.append(file_text)
                            file_count += 1
                        except Exception:
                            continue
                    
                    # Limit processing to avoid excessive memory usage
                    if file_count >= 20:
                        break
            
            full_text = '\n'.join(text_parts)
            if not full_text.strip():
                full_text = f"ZIP archive containing {len(zip_ref.filelist)} files (no extractable text found)"
            
            pages = self._split_into_pages(full_text)
            
            return {
                "text": full_text,
                "pages": pages,
                "metadata": {
                    "files_processed": file_count,
                    "total_files": len(zip_ref.filelist),
                    "page_count": len(pages),
                    "extractor": "zipfile"
                }
            }
            
        except Exception as e:
            logger.error(f"Failed to extract text from ZIP: {str(e)}")
            return {"text": "", "pages": [], "metadata": {"error": str(e)}}
    
    def _split_into_pages(self, text: str, chars_per_page: int = 2500) -> List[Dict[str, Any]]:
        """Split text into pages based on character count or natural breaks."""
        if not text.strip():
            return []
        
        pages = []
        lines = text.split('\n')
        current_page = []
        current_length = 0
        page_num = 1
        
        for line in lines:
            line_length = len(line) + 1  # +1 for newline
            
            # If adding this line would exceed the limit, finalize current page
            if current_length + line_length > chars_per_page and current_page:
                pages.append({
                    "page": page_num,
                    "text": '\n'.join(current_page)
                })
                current_page = [line]
                current_length = line_length
                page_num += 1
            else:
                current_page.append(line)
                current_length += line_length
        
        # Add final page
        if current_page:
            pages.append({
                "page": page_num,
                "text": '\n'.join(current_page)
            })
        
        return pages


class TextChunker:
    """Split text into chunks with overlap for embedding."""
    
    def __init__(self, chunk_size: int = 1000, overlap: int = 200):
        self.chunk_size = chunk_size
        self.overlap = overlap
    
    def chunk_text(self, text: str, page_info: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        """Split text into chunks with overlap."""
        if not text.strip():
            return []
        
        # Simple token approximation (1 token ≈ 4 characters)
        chunk_chars = self.chunk_size * 4
        overlap_chars = self.overlap * 4
        
        chunks = []
        start = 0
        chunk_num = 1
        
        while start < len(text):
            # Calculate end position
            end = min(start + chunk_chars, len(text))
            
            # Try to break at sentence or paragraph boundaries
            chunk_text = text[start:end]
            
            # If not at the end of text, try to break at sentence boundary
            if end < len(text):
                # Look for sentence endings in the last 200 characters
                sentence_break = chunk_text.rfind('. ')
                if sentence_break > len(chunk_text) - 200:
                    end = start + sentence_break + 1
                    chunk_text = text[start:end]
            
            # Determine page range for this chunk
            start_page, end_page = self._find_page_range(start, end, page_info, text)
            
            chunks.append({
                "chunk_number": chunk_num,
                "text": chunk_text.strip(),
                "start_char": start,
                "end_char": end,
                "start_page": start_page,
                "end_page": end_page,
                "token_count": len(chunk_text) // 4  # Rough estimate
            })
            
            # Move to next chunk with overlap
            start = max(start + chunk_chars - overlap_chars, end)
            chunk_num += 1
            
            # Prevent infinite loop
            if start >= len(text):
                break
        
        return chunks
    
    def _find_page_range(self, start_char: int, end_char: int, page_info: List[Dict[str, Any]], full_text: str) -> tuple:
        """Find which pages this chunk spans."""
        if not page_info:
            return 1, 1
        
        start_page = 1
        end_page = 1
        current_pos = 0
        
        for page in page_info:
            page_text = page.get('text', '')
            page_end = current_pos + len(page_text)
            
            # Check if chunk starts in this page
            if current_pos <= start_char < page_end:
                start_page = page.get('page', 1)
            
            # Check if chunk ends in this page
            if current_pos <= end_char <= page_end:
                end_page = page.get('page', 1)
                break
            
            current_pos = page_end + 2  # +2 for page separators
        
        return start_page, end_page


class DocumentProcessor:
    """Main document processing orchestrator."""
    
    def __init__(self):
        self.extractor = TextExtractor()
        self.chunker = TextChunker()
        self.embedder = EmbedderClient()
        self.redis_client = None
        self.engine = None
    
    async def initialize(self):
        """Initialize connections."""
        self.redis_client = redis.from_url(settings.redis_url)
        self.engine = create_async_engine(settings.database_url)
    
    async def process_document(self, task_data: Dict[str, Any]) -> None:
        """Process a document: extract text, chunk, and embed."""
        document_id = UUID(task_data['document_id'])
        version_number = task_data['version_number']
        storage_path = task_data['storage_path']
        content_type = task_data['content_type']
        
        logger.info(f"Processing document {document_id}, version {version_number}")
        
        try:
            # Get file from storage
            storage = StorageService()
            file_data = await storage.retrieve_file(storage_path)
            
            # Extract text
            extraction_result = await self.extractor.extract_text(
                file_data, content_type, f"doc_{document_id}"
            )
            
            if not extraction_result['text'].strip():
                logger.warning(f"No text extracted from document {document_id}")
                return
            
            # Chunk text
            chunks = self.chunker.chunk_text(
                extraction_result['text'],
                extraction_result['pages']
            )
            
            if not chunks:
                logger.warning(f"No chunks created for document {document_id}")
                return
            
            # Process chunks and create fragments
            async with AsyncSession(self.engine) as session:
                for chunk in chunks:
                    await self._create_fragment(
                        session,
                        document_id,
                        version_number,
                        chunk
                    )
                
                await session.commit()
            
            logger.info(f"Successfully processed document {document_id}: {len(chunks)} fragments created")
            
        except Exception as e:
            logger.error(f"Failed to process document {document_id}: {str(e)}")
            raise
    
    async def _create_fragment(
        self,
        session: AsyncSession,
        document_id: UUID,
        version_number: int,
        chunk: Dict[str, Any]
    ) -> None:
        """Create fragment record with embedding."""
        # Get embedding for chunk text
        embedding = await self.embedder.get_embedding(chunk['text'])
        
        # Create fragment
        fragment_id = UUID(str(uuid4()))
        
        # Prepare embedding for PostgreSQL vector type
        embedding_str = f"[{','.join(map(str, embedding))}]" if embedding else None
        
        stmt = text("""
            INSERT INTO fragments (
                id, document_id, version_number, content, content_preview,
                start_page, end_page, fragment_type, language, confidence_score,
                classification_level, metadata, tags, embedding, created_at
            ) VALUES (
                :id, :document_id, :version_number, :content, :content_preview,
                :start_page, :end_page, :fragment_type, :language, :confidence_score,
                :classification_level, :metadata, :tags, :embedding::vector, :created_at
            )
        """)
        
        await session.execute(stmt, {
            "id": fragment_id,
            "document_id": document_id,
            "version_number": version_number,
            "content": chunk['text'],
            "content_preview": chunk['text'][:200] + "..." if len(chunk['text']) > 200 else chunk['text'],
            "start_page": chunk['start_page'],
            "end_page": chunk['end_page'],
            "fragment_type": "text",
            "language": "en",
            "confidence_score": 0.95,
            "classification_level": "restricted",
            "metadata": {
                "chunk_number": chunk['chunk_number'],
                "token_count": chunk['token_count'],
                "start_char": chunk['start_char'],
                "end_char": chunk['end_char']
            },
            "tags": [],
            "embedding": embedding_str,
            "created_at": datetime.utcnow()
        })


async def main():
    """Main worker loop."""
    processor = DocumentProcessor()
    await processor.initialize()
    
    logger.info("Document processor worker started")
    
    while True:
        try:
            # Wait for tasks from Redis queue
            task_data = await processor.redis_client.brpop("ingest_queue", timeout=10)
            
            if task_data:
                # Parse task data
                queue_name, task_json = task_data
                task = json.loads(task_json)
                
                # Process document
                await processor.process_document(task)
            
        except KeyboardInterrupt:
            logger.info("Shutdown signal received")
            break
        except Exception as e:
            logger.error(f"Worker error: {str(e)}")
            # Continue processing other tasks
            continue
    
    # Cleanup
    if processor.redis_client:
        await processor.redis_client.close()
    if processor.engine:
        await processor.engine.dispose()


if __name__ == "__main__":
    # Import missing datetime
    from datetime import datetime
    from uuid import uuid4
    
    asyncio.run(main())